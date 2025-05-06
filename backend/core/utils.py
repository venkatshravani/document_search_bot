from PyPDF2 import PdfReader
import docx
import openpyxl
from pptx import Presentation

def extract_pdf_text(file):
    reader = PdfReader(file)
    texts = []
    for page in reader.pages:
        texts.append(page.extract_text() or '')
    return ' '.join(texts)

def extract_word_text(file):
    doc = docx.Document(file)
    texts = [para.text for para in doc.paragraphs if para.text]
    return ' '.join(texts)

def extract_excel_text(file):
    wb = openpyxl.load_workbook(file)
    texts = []
    for sheet in wb.sheetnames:
        sheet_data = wb[sheet]
        for row in sheet_data.iter_rows():
            texts.extend([str(cell.value) if cell.value is not None else '' for cell in row])
    return ' '.join(texts)

def extract_pptx_text(file):
    prs = Presentation(file)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
    return ' '.join(texts)
